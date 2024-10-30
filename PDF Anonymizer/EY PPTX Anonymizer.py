# Databricks notebook source
# MAGIC %pip
# MAGIC install python-pptx
# MAGIC %pip install presidio-analyzer
# MAGIC %pip install presidio-anonymizer

# COMMAND ----------

from pptx import Presentation
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine
import os
import uuid

pptx_path = "/Volumes/yash_gupta_catalog/ey_dbs/dbs_pptx/Project Helios DRAFT Report - 05 30 17.pptx"
anonymized_pptx_path = "/Volumes/yash_gupta_catalog/ey_dbs/dbs_pptx/Project Helios DRAFT Report - 05 30 17 Anonymized.pptx"

# COMMAND ----------

import os
import uuid
from pptx import Presentation
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine

analyzer = AnalyzerEngine()
anonymizer = AnonymizerEngine()

if os.path.exists(pptx_path):
    prs = Presentation(pptx_path)
    slides_content = []
    
    entity_mapping = []
    raw_documents = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                results = analyzer.analyze(
                    text=text, 
                    entities=["PERSON", "ORGANIZATION"], 
                    language="en"
                )
                anonymized_text = anonymizer.anonymize(
                    text=text, 
                    analyzer_results=results
                ).text
                shape.text = anonymized_text
                
                for result in results:
                    entity_text = text[result.start:result.end]
                    entity_mapping.append(
                        (result.entity_type, entity_text, str(uuid.uuid4()))
                    )
                
                raw_documents.append(
                    ("Project Helios DRAFT Report - 05 30 17.pptx", slide_num, text)
                )
    
    # Save to a temporary local path
    temp_local_path = "/tmp/anonymized_presentation.pptx"
    prs.save(temp_local_path)
    
    # Move the file to the desired UC volume location
    dbutils.fs.mv("file:/tmp/anonymized_presentation.pptx", anonymized_pptx_path)
    
    entity_mapping_df = spark.createDataFrame(
        entity_mapping, 
        ["entity_type", "entity_text", "entity_unique_id"]
    )
    entity_mapping_df.write.mode("overwrite").saveAsTable(
        "yash_gupta_catalog.ey_dbs.entity_mapping_table"
    )
    
    raw_documents_df = spark.createDataFrame(
        raw_documents, 
        ["doc_name", "page_num", "text"]
    )
    raw_documents_df.write.mode("overwrite").saveAsTable("yash_gupta_catalog.ey_dbs.raw_documents_table")
else:
    print(f"File not found: {pptx_path}")

# COMMAND ----------

display(entity_mapping_df)
display(raw_documents_df)

